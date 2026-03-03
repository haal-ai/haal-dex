"""Output generator for rendering, validating, and exporting documents.

Uses Jinja2 for template rendering, and delegates to WeasyPrint (PDF),
python-docx (DOCX) for export.  Encryption is delegated
to :class:`EncryptionService` when the template has encryption settings.
"""

from __future__ import annotations

import os
import re
import sys
import uuid
from datetime import datetime, timezone
from pathlib import Path

from jinja2 import Environment, FileSystemLoader, TemplateNotFound

from app.config import Settings, get_settings
from app.models.templates import (
    DocumentMetadata,
    RenderedDocument,
    Template,
    ValidationResult,
    ValidationRule,
)
from app.services.encryption_service import EncryptionService
from app.services.template_registry import TemplateRegistry


class OutputGenerator:
    """Renders, validates, and exports documents from templates."""

    def __init__(
        self,
        registry: TemplateRegistry,
        encryption_service: EncryptionService | None = None,
        settings: Settings | None = None,
    ) -> None:
        self._registry = registry
        self._encryption = encryption_service or EncryptionService()
        self._settings = settings or get_settings()

    # ------------------------------------------------------------------
    # render
    # ------------------------------------------------------------------

    async def render(
        self,
        template_id: str,
        data: dict,
        metadata: DocumentMetadata,
        session_id: str = "",
    ) -> RenderedDocument:
        """Render a document by applying a Jinja2 template to *data*.

        The template is looked up in the :class:`TemplateRegistry` and its
        ``jinja2_template_path`` is loaded from the filesystem.  The
        *metadata* fields are injected into the template context alongside
        *data*.

        Returns a :class:`RenderedDocument` whose ``content`` is the
        rendered bytes (UTF-8 encoded).
        """
        template = self._registry.get_template(template_id)

        rendered_text = self._render_jinja2(template, data, metadata)

        content = rendered_text.encode("utf-8")

        # Encrypt if the template has encryption settings enabled.
        if template.encryption_settings and template.encryption_settings.enabled:
            content = self._encryption.encrypt(content, template.encryption_settings)

        validation_result = self._run_validation(rendered_text, template)

        return RenderedDocument(
            id=str(uuid.uuid4()),
            session_id=session_id,
            template_id=template_id,
            format=template.format,
            content=content,
            metadata=metadata,
            validation_result=validation_result,
        )

    # ------------------------------------------------------------------
    # validate
    # ------------------------------------------------------------------

    async def validate(
        self, document: RenderedDocument, template: Template
    ) -> ValidationResult:
        """Check *document* against *template* validation rules.

        Returns a list of human-readable violation strings.  An empty list
        means the document is valid.
        """
        text = document.content.decode("utf-8", errors="replace")
        return self._run_validation(text, template)

    # ------------------------------------------------------------------
    # export
    # ------------------------------------------------------------------

    async def export(self, document: RenderedDocument, fmt: str) -> bytes:
        """Export *document* to the requested *fmt*.

        Supported formats: ``"pdf"``, ``"docx"``, ``"md"``,
        ``"html"``, ``"pptx"``.
        """
        text = document.content.decode("utf-8", errors="replace")

        if document.format == "md":
            if fmt == "pdf":
                html_text = self._markdown_to_html(text)
                return self._export_pdf(html_text)
            elif fmt == "html":
                html_text = self._markdown_to_html(text)
                return html_text.encode("utf-8")
            elif fmt == "docx":
                return self._export_docx_from_markdown(text, document.metadata)

        if fmt == "pdf":
            return self._export_pdf(text)
        elif fmt == "docx":
            return self._export_docx(text, document.metadata)
        elif fmt in ("md", "markdown"):
            return text.encode("utf-8")
        elif fmt == "html":
            return text.encode("utf-8")
        elif fmt == "pptx":
            return self._export_pptx(text)
        else:
            raise ValueError(f"Unsupported export format: {fmt!r}")

    # ------------------------------------------------------------------
    # Internal – Jinja2 rendering
    # ------------------------------------------------------------------

    def _render_jinja2(
        self, template: Template, data: dict, metadata: DocumentMetadata
    ) -> str:
        """Load the Jinja2 template file and render it."""
        tpl_path = Path(template.jinja2_template_path)

        env = Environment(
            loader=FileSystemLoader(str(tpl_path.parent)),
            autoescape=False,
        )

        try:
            jinja_tpl = env.get_template(tpl_path.name)
        except TemplateNotFound:
            raise FileNotFoundError(
                f"Jinja2 template not found: {template.jinja2_template_path}"
            )

        context = {
            **data,
            "metadata": {
                "author": metadata.author,
                "date": metadata.date.isoformat() if isinstance(metadata.date, datetime) else str(metadata.date),
                "version": metadata.version,
                "classification": metadata.classification,
            },
        }

        return jinja_tpl.render(context)

    @staticmethod
    def _markdown_to_html(markdown_text: str) -> str:
        try:
            import markdown as md  # type: ignore[import-untyped]
        except Exception as exc:
            raise RuntimeError(
                "Markdown rendering requires the 'Markdown' Python package. "
                "Install backend dependencies and restart the backend."
            ) from exc

        return md.markdown(
            markdown_text,
            extensions=["extra"],
            output_format="html5",
        )

    # ------------------------------------------------------------------
    # Internal – validation
    # ------------------------------------------------------------------

    def _run_validation(self, text: str, template: Template) -> ValidationResult:
        """Run all validation rules against *text*."""
        violations: ValidationResult = []
        for rule in template.validation_rules:
            violation = self._check_rule(text, rule)
            if violation:
                violations.append(violation)
        return violations

    @staticmethod
    def _check_rule(text: str, rule: ValidationRule) -> str | None:
        """Return a violation message if *rule* is violated, else ``None``."""
        if rule.rule_type == "required":
            # The field value must appear in the text.
            if rule.field not in text:
                return f"Required field '{rule.field}' is missing from the document"

        elif rule.rule_type == "regex":
            pattern = rule.parameters.get("pattern", "")
            if pattern and not re.search(pattern, text):
                return (
                    f"Field '{rule.field}' does not match required pattern "
                    f"'{pattern}'"
                )

        elif rule.rule_type == "format":
            expected = rule.parameters.get("expected", "")
            if expected and expected not in text:
                return (
                    f"Field '{rule.field}' does not match expected format "
                    f"'{expected}'"
                )

        elif rule.rule_type == "cross_reference":
            ref_field = rule.parameters.get("reference_field", "")
            if ref_field and ref_field not in text:
                return (
                    f"Cross-reference for '{rule.field}' failed: "
                    f"referenced field '{ref_field}' not found"
                )

        return None

    # ------------------------------------------------------------------
    # Internal – export helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _export_pdf(html_text: str) -> bytes:
        """Convert HTML text to PDF via WeasyPrint."""
        if sys.platform.startswith("win"):
            dll_dirs_env = os.getenv("INTENT_WEASYPRINT_DLL_DIRS")
            dll_dirs = (
                [p for p in dll_dirs_env.split(";") if p.strip()]
                if dll_dirs_env
                else [r"C:\\msys64\\mingw64\\bin", r"C:\\msys64\\usr\\bin"]
            )
            for p in dll_dirs:
                try:
                    if Path(p).exists():
                        os.add_dll_directory(p)
                except Exception:
                    pass

        try:
            from weasyprint import HTML  # type: ignore[import-untyped]
        except Exception as exc:
            raise RuntimeError(
                "PDF export requires WeasyPrint and its native dependencies. "
                "On Windows, install a GTK/Pango/Cairo runtime (e.g., MSYS2 mingw-w64 or a GTK3 runtime), "
                "then restart the backend. See: https://doc.courtbouillon.org/weasyprint/stable/first_steps.html#installation"
            ) from exc

        try:
            return HTML(string=html_text).write_pdf()
        except Exception as exc:
            raise RuntimeError(
                "WeasyPrint failed to generate the PDF. This is usually caused by missing native libraries "
                "(Cairo/Pango/GDK-PixBuf) on Windows. See: https://doc.courtbouillon.org/weasyprint/stable/first_steps.html#troubleshooting"
            ) from exc

    @staticmethod
    def _export_docx(text: str, metadata: DocumentMetadata) -> bytes:
        """Build a DOCX document from plain text with metadata."""
        from io import BytesIO

        from docx import Document  # type: ignore[import-untyped]

        doc = Document()
        doc.core_properties.author = metadata.author
        doc.core_properties.title = f"v{metadata.version}"
        doc.core_properties.comments = metadata.classification

        for line in text.split("\n"):
            doc.add_paragraph(line)

        buf = BytesIO()
        doc.save(buf)
        return buf.getvalue()

    @staticmethod
    def _export_docx_from_markdown(markdown_text: str, metadata: DocumentMetadata) -> bytes:
        from io import BytesIO

        from docx import Document  # type: ignore[import-untyped]

        doc = Document()
        doc.core_properties.author = metadata.author
        doc.core_properties.title = f"v{metadata.version}"
        doc.core_properties.comments = metadata.classification

        for raw_line in markdown_text.splitlines():
            line = raw_line.rstrip("\r\n")
            stripped = line.strip()
            if not stripped:
                doc.add_paragraph("")
                continue

            if stripped.startswith("### "):
                doc.add_heading(stripped[4:].strip(), level=3)
            elif stripped.startswith("## "):
                doc.add_heading(stripped[3:].strip(), level=2)
            elif stripped.startswith("# "):
                doc.add_heading(stripped[2:].strip(), level=1)
            elif stripped.startswith("- ") or stripped.startswith("* "):
                doc.add_paragraph(stripped[2:].strip(), style="List Bullet")
            else:
                doc.add_paragraph(stripped)

        buf = BytesIO()
        doc.save(buf)
        return buf.getvalue()

    @staticmethod
    def _export_pptx(text: str) -> bytes:
        """Convert a simple markdown-like slide outline to PPTX.

        Slides are separated by lines containing only '---'.
        The first non-empty line of each slide becomes the title.
        Remaining lines become the body (joined with newlines).
        """
        from io import BytesIO

        from pptx import Presentation  # type: ignore[import-untyped]

        prs = Presentation()
        # Title and content layout is usually index 1
        layout = prs.slide_layouts[1]

        chunks = [c.strip() for c in re.split(r"\r?\n---\r?\n", text)]
        for chunk in [c for c in chunks if c.strip()]:
            lines = [ln.strip() for ln in chunk.splitlines() if ln.strip()]
            if not lines:
                continue
            title = lines[0].lstrip("# ")
            body = "\n".join(lines[1:]) if len(lines) > 1 else ""

            slide = prs.slides.add_slide(layout)
            if slide.shapes.title is not None:
                slide.shapes.title.text = title
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                tf.text = body

        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()
